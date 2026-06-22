# -*- coding: utf-8 -*-
"""Сборка обновлённой презентации: +4 слайда в стиле деки, перенумерация футеров."""
import os, re, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = r"C:\Users\krealix\Desktop\Презентация ВКР — Сульдин1.pptx"
DST = r"C:\Users\krealix\Desktop\Презентация ВКР — Сульдин (обновлённая).pptx"
AST = r"C:\Users\krealix\Desktop\Новая папка (3)\_assets"

NAVY="0E2A47"; TEAL="0E9AA7"; TEAL2="27C2B6"; GREEN="2BA84A"; RED="D9534F"
MUTED="5B6B7B"; INK="1E2A38"; LINE="D4DEE8"; WHITE="FFFFFF"
def C(h): return RGBColor.from_string(h)

prs = Presentation(SRC)
LAYOUT = prs.slides[7].slide_layout
EMU_IN = 914400

def blank():
    s = prs.slides.add_slide(LAYOUT)
    for shp in list(s.shapes):
        shp._element.getparent().remove(shp._element)
    return s

def _set_text(tf, parts, anchor=MSO_ANCHOR.TOP):
    """parts: list of (text, name, size, bold, color, align). First fills p0."""
    tf.word_wrap = True
    try: tf.vertical_anchor = anchor
    except Exception: pass
    for i,(txt,name,size,bold,color,align) in enumerate(parts):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        if align is not None: p.alignment = align
        p.space_before = Pt(0); p.space_after = Pt(2); p.line_spacing = 1.05
        r = p.add_run(); r.text = txt
        f = r.font; f.name = name; f.size = Pt(size); f.bold = bold; f.color.rgb = C(color)

def textbox(s, x,y,w,h, parts, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf = tb.text_frame
    tf.margin_left=Pt(0); tf.margin_right=Pt(0); tf.margin_top=Pt(0); tf.margin_bottom=Pt(0)
    _set_text(tf, parts, anchor)
    return tb

def eyebrow(s,t): textbox(s,0.7,0.42,11.93,0.32,[(t,"Trebuchet MS",12.5,True,TEAL,None)])
def title(s,t):   textbox(s,0.7,0.74,11.93,0.85,[(t,"Trebuchet MS",28,True,NAVY,None)])
def footer(s):
    textbox(s,11.63,7.0,1.1,0.3,[("00 / 00","Calibri",10,False,MUTED,PP_ALIGN.RIGHT)])

def card(s, n, x, y, w, ttl, body, badge=TEAL, tw=None):
    # badge
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),Inches(y),Inches(0.46),Inches(0.46))
    b.adjustments[0]=0.18
    b.fill.solid(); b.fill.fore_color.rgb=C(badge); b.line.fill.background()
    b.shadow.inherit=False
    btf=b.text_frame; btf.margin_left=Pt(0); btf.margin_right=Pt(0); btf.margin_top=Pt(0); btf.margin_bottom=Pt(0)
    btf.word_wrap=False
    p=btf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(n); r.font.name="Trebuchet MS"; r.font.size=Pt(17); r.font.bold=True; r.font.color.rgb=C(WHITE)
    try: btf.vertical_anchor=MSO_ANCHOR.MIDDLE
    except Exception: pass
    # text
    width = tw if tw else (w-0.62)
    textbox(s, x+0.62, y-0.06, width, 1.1,
            [(ttl,"Calibri",14,True,NAVY,None),(body,"Calibri",12.5,False,INK,None)])

def image_framed(s, path, x, y, w, h, pad=0.14):
    fr = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x-pad),Inches(y-pad),Inches(w+2*pad),Inches(h+2*pad))
    fr.adjustments[0]=0.035
    fr.fill.solid(); fr.fill.fore_color.rgb=C(WHITE)
    fr.line.color.rgb=C(LINE); fr.line.width=Pt(1.0); fr.shadow.inherit=False
    s.shapes.add_picture(path, Inches(x),Inches(y),Inches(w),Inches(h))

def img_size(path):
    from PIL import Image
    with Image.open(path) as im: return im.size

def clear_slide(s):
    for shp in list(s.shapes):
        shp._element.getparent().remove(shp._element)

def section_header(s, x, y, w, txt):
    textbox(s, x, y, w, 0.4, [(txt, "Trebuchet MS", 15, True, TEAL, None)])

def bullet(s, x, y, w, txt):
    d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y+0.07), Inches(0.15), Inches(0.15))
    d.fill.solid(); d.fill.fore_color.rgb=C(TEAL); d.line.fill.background(); d.shadow.inherit=False
    textbox(s, x+0.4, y-0.04, w, 0.8, [(txt, "Calibri", 12.5, False, INK, None)])

def tech_row(s, x, y, label, desc, pillw=2.5, descw=3.55):
    p = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(pillw), Inches(0.44))
    p.adjustments[0]=0.5
    p.fill.solid(); p.fill.fore_color.rgb=C("E1F1F2")
    p.line.color.rgb=C(TEAL); p.line.width=Pt(1.0); p.shadow.inherit=False
    tf=p.text_frame; tf.word_wrap=False
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); tf.margin_top=Pt(0); tf.margin_bottom=Pt(0)
    pp=tf.paragraphs[0]; pp.alignment=PP_ALIGN.CENTER
    r=pp.add_run(); r.text=label; r.font.name="Calibri"; r.font.size=Pt(11.5); r.font.bold=True; r.font.color.rgb=C("0E6E78")
    try: tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    except Exception: pass
    parts=[(ln,"Calibri",11.5,False,INK,None) for ln in desc.split("\n")]
    textbox(s, x+pillw+0.25, y-0.04, descw, 0.9, parts)

# ---------------- Слайд: Генератор STOWN ----------------
s = blank(); eyebrow(s,"ГЛАВА 2 · ПРОЕКТНАЯ ЧАСТЬ"); title(s,"Генератор метки STOWN"); footer(s)
cards=[("Пакет 10 байт","Команда открытия + 7-байтный идентификатор носителя + номер замка (uint16, BE). Фиксированная длина — простой и надёжный разбор."),
       ("Идентификатор носителя","Стабильный 7-байтный токен (56 бит); поддержка MAC, UUID и телефона/IMEI в BCD-формате."),
       ("Динамический rolling-code","ID = HMAC-SHA256(секрет, шаг 30 с): код меняется каждые 30 с — защита от клонирования и ретрансляции."),
       ("Обёртки рекламы","Manufacturer Data · Service Data · iBeacon; непрерывное вещание (вкладка «Метка», flutter_ble_peripheral).")]
yy=2.0
for i,(t,b) in enumerate(cards):
    card(s,i+1,0.7,yy,5.3,t,b,tw=4.55); yy+=1.18
# packet figure (landscape) right
p="%s/fig_stown_packet.png"%AST; W,H=img_size(p); w=7.05; h=w*H/W
image_framed(s,p,5.85,2.45,w,h)

# (Канал по звонку теперь встроен в архитектуру слайда 7 — отдельный слайд не создаётся)

# ---------------- Слайд: RSSI и расстояние ----------------
s = blank(); eyebrow(s,"ГЛАВА 3 · РЕАЛИЗАЦИЯ И ТЕСТИРОВАНИЕ"); title(s,"Натурный эксперимент: RSSI и расстояние"); footer(s)
p="%s/fig_dist_rssi.png"%AST; W,H=img_size(p); w=7.9; h=w*H/W
image_framed(s,p,0.55,1.95,w,h)
cards=[("Сигнал убывает с расстоянием","Средний RSSI снижается с −65 дБм на 1 м до −90 дБм на 15 м — основа оценки близости."),
       ("Плато 2–5 м","На 2–5 м уровень почти неизменен (≈ −70 дБм): мгновенный порог не различает эти расстояния."),
       ("Вывод","В зоне неоднозначности один порог ненадёжен → решение по траектории и гистерезису зон сигнала.")]
yy=2.05
for i,(t,b) in enumerate(cards):
    card(s,i+1,8.75,yy,4.1,t,b,badge=TEAL,tw=3.5); yy+=1.5

# ---------------- Слайд: Гистерезис на реальных данных ----------------
s = blank(); eyebrow(s,"ГЛАВА 3 · РЕАЛИЗАЦИЯ И ТЕСТИРОВАНИЕ"); title(s,"Гистерезис на реальных данных"); footer(s)
p="%s/fig_hyst_real.png"%AST; W,H=img_size(p); h=5.55; w=h*W/H
image_framed(s,p,0.5,1.55,w,h)
cards=[("Реальный проход носителя","Журнал сканера: STOWN-метка, ~4 изм./с; пороги «близко» −65 дБм, «далеко» −75 дБм."),
       ("Взвод → удержание → доступ","Метка проходит «далеко» (взвод B), затем удерживается «близко» (A > Y) → доступ выдан на t ≈ 43 с."),
       ("Защёлка гистерезиса","Повторный подход без нового взвода доступа не открывает — защита от ложных и повторных срабатываний.")]
yy=2.0
for i,(t,b) in enumerate(cards):
    card(s,i+1,8.55,yy,4.3,t,b,badge=TEAL,tw=3.7); yy+=1.5

# ---------------- Слайд 5: «Методы и средства» — инфографика с логотипами ----------------
s5 = prs.slides[4]
clear_slide(s5)
eyebrow(s5, "04 · ВВЕДЕНИЕ"); title(s5, "Методы и средства"); footer(s5)
p5 = "%s/fig_slide5.png" % AST
W5, H5 = img_size(p5)
H_TARGET = 4.6
w5 = H_TARGET * W5 / H5
x5 = (13.333 - w5) / 2
s5.shapes.add_picture(p5, Inches(x5), Inches(1.62), Inches(w5), Inches(H_TARGET))
textbox(s5, x5 + 0.1, 1.62 + H_TARGET + 0.16, w5 - 0.2, 0.45,
        [("Единая кодовая база (Flutter): весь цикл в одном приложении — генерация метки, "
          "приём и анализ сигнала, решение и отправка команды замку.", "Calibri", 11, False, MUTED, None)])

# ---------------- Слайд 7: встраиваем канал по звонку в архитектуру ----------------
def set_text_keep_fmt(shape, txt):
    tf = shape.text_frame
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    runs = p0.runs
    if runs:
        runs[0].text = txt
        for r in list(runs[1:]):
            r._r.getparent().remove(r._r)
    else:
        p0.add_run().text = txt

slide7 = prs.slides[6]  # ГЛАВА 2 · Архитектура системы
pic = next((sh for sh in slide7.shapes if sh.shape_type == 13), None)
if pic is not None:
    L, T, Wd, Ht = pic.left, pic.top, pic.width, pic.height
    pic._element.getparent().remove(pic._element)
    slide7.shapes.add_picture("%s/fig_arch_dual.png" % AST, L, T, Wd, Ht)
for sh in slide7.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith("Модульность"):
        set_text_keep_fmt(sh, "Два канала доступа: BLE с анализом траектории (основной) "
                               "и доступ по входящему звонку (резерв/гость).")

# ---------------- Слайд 8: красивая блок-схема алгоритма ----------------
slide8 = prs.slides[7]
pic8 = next((sh for sh in slide8.shapes if sh.shape_type == 13), None)
if pic8 is not None:
    bx, by = Emu(pic8.left).inches, Emu(pic8.top).inches
    bw, bh = Emu(pic8.width).inches, Emu(pic8.height).inches
    pic8._element.getparent().remove(pic8._element)
    p8 = "%s/fig_algo.png" % AST
    W8, H8 = img_size(p8); ar = W8 / H8
    if bw / bh > ar:           # коробка шире картинки → вписать по высоте
        nh = bh; nw = nh * ar; nx = bx + (bw - nw) / 2; ny = by
    else:                      # вписать по ширине
        nw = bw; nh = nw / ar; nx = bx; ny = by + (bh - nh) / 2
    slide8.shapes.add_picture(p8, Inches(nx), Inches(ny), Inches(nw), Inches(nh))

# ---------------- Перестановка порядка слайдов ----------------
# текущий: 0..9 (ориг) + 10 STOWN, 11 DIST, 12 HYST
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
order = [0,1,2,3,4,5,6,7, 10, 8, 11,12, 9]
for i in order:
    sldIdLst.append(ids[i])

# ---------------- Перенумерация футеров ----------------
total = len(prs.slides.__iter__.__self__._sldIdLst)  # = number of sldId
total = len(list(prs.slides))
pat = re.compile(r'^\s*\d{1,2}\s*/\s*\d{1,2}\s*$')
for idx, sl in enumerate(prs.slides):
    for sh in sl.shapes:
        if sh.has_text_frame and pat.match(sh.text_frame.text or ""):
            # заменить текст, сохранив форматирование первого run
            tf = sh.text_frame
            run = None
            for para in tf.paragraphs:
                for r in para.runs:
                    run = r; break
                if run: break
            if run is not None:
                run.text = "%02d / %d" % (idx+1, total)

prs.save(DST)
print("saved:", DST, "total slides:", total)
