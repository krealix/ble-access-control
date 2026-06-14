"""Сборка презентации ВКР (.pptx) из готовых рисунков.

Светлый фон, тёмный текст, без анимаций (по требованиям методички).
Титульный слайд + 9 содержательных + «Спасибо за внимание».
"""
import os
import struct

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

D = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(D, "figures")
OUT = os.path.join(D, "ВКР_презентация.pptx")

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x14, 0x21, 0x36)
ACCENT = RGBColor(0x2D, 0x8C, 0xFF)
DARK = RGBColor(0x22, 0x2A, 0x33)
MUTED = RGBColor(0x55, 0x5F, 0x6D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = 13.333


def png_size(path):
    with open(path, "rb") as f:
        f.read(16)
        return struct.unpack(">II", f.read(8))


def blank():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = WHITE
    return s


def title(slide, text, size=30):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = "Arial"
    r.font.color.rgb = NAVY
    line = slide.shapes.add_shape(1, Inches(0.65), Inches(1.5), Inches(3.2), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def bullets(slide, items, left=0.8, top=1.85, width=11.8, height=5.0, size=20):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        sub = it.startswith("\t")
        text = it.strip()
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = ("–  " if sub else "•  ") + text
        r.font.size = Pt(size - 3 if sub else size)
        r.font.name = "Arial"
        r.font.color.rgb = MUTED if sub else DARK
        p.space_after = Pt(8)
        p.level = 1 if sub else 0


def image_fit(slide, name, top=1.85, max_w=11.6, max_h=5.0, caption=None):
    path = os.path.join(FIG, name)
    w, h = png_size(path)
    ratio = w / h
    width = max_w
    height = width / ratio
    if height > max_h:
        height = max_h
        width = height * ratio
    left = (SW - width) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    if caption:
        cb = slide.shapes.add_textbox(Inches(0.6), Inches(top + height + 0.05),
                                      Inches(12.1), Inches(0.5))
        cp = cb.text_frame.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        rr = cp.add_run()
        rr.text = caption
        rr.font.size = Pt(13)
        rr.font.name = "Arial"
        rr.font.italic = True
        rr.font.color.rgb = MUTED


# --- Слайд 1: титульный ---
s = blank()
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.2))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = "Разработка системы контроля удалённого доступа на основе анализа траектории изменения сигнала BLE-меток"
r.font.size = Pt(30)
r.font.bold = True
r.font.name = "Arial"
r.font.color.rgb = NAVY
info = s.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.6)).text_frame
info.word_wrap = True
for i, line_txt in enumerate([
    "Выпускная квалификационная работа (магистерская диссертация)",
    "Направление 01.04.02 — Прикладная математика и информатика",
    "",
    "Выполнил(а): ________________________      Группа: __________",
    "Научный руководитель: ________________________",
    "",
    "Ханты-Мансийск, 2026",
]):
    pp = info.paragraphs[0] if i == 0 else info.add_paragraph()
    rr = pp.add_run()
    rr.text = line_txt
    rr.font.size = Pt(16)
    rr.font.name = "Arial"
    rr.font.color.rgb = DARK if i == 0 else MUTED

# --- Слайд 2: актуальность ---
s = blank()
title(s, "Актуальность и постановка проблемы")
bullets(s, [
    "Бесконтактные системы контроля доступа (СКУД) на базе BLE — доступ «по присутствию» носителя метки",
    "Близость носителя оценивается по уровню сигнала RSSI",
    "RSSI сильно зашумлён: переотражения, экранирование, помехи",
    "Решение по порогу RSSI ненадёжно: ложные срабатывания, нет учёта направления, уязвимость к подмене",
    "Решение: анализ ТРАЕКТОРИИ изменения сигнала во времени",
])

# --- Слайд 3: цель и задачи ---
s = blank()
title(s, "Объект, предмет, цель и задачи")
bullets(s, [
    "Объект: системы контроля удалённого доступа на основе BLE-меток",
    "Предмет: методы и алгоритмы анализа траектории RSSI для решения о доступе",
    "Цель: разработать систему, принимающую решение по траектории сигнала, и подтвердить работоспособность",
    "Задачи:",
    "\tобзор и выбор метода",
    "\tархитектура системы",
    "\tалгоритм анализа траектории",
    "\tпрограммная реализация",
    "\tтестирование",
])

# --- Слайд 4: RSSI и расстояние ---
s = blank()
title(s, "Сигнал RSSI и оценка расстояния")
bullets(s, [
    "Модель затухания: RSSI = A − 10·n·lg(d); расстояние d = 10^((A−RSSI)/(10n))",
    "Сигнал зашумлён и нестабилен → нужна фильтрация и анализ динамики",
], top=1.7, height=1.4, size=18)
image_fit(s, "fig_1_1_rssi_distance.png", top=3.2, max_h=3.9,
          caption="Зависимость RSSI от расстояния при разных n")

# --- Слайд 5: архитектура ---
s = blank()
title(s, "Архитектура системы")
image_fit(s, "fig_2_1_architecture.png", top=2.3, max_h=3.6,
          caption="метка → сканер → анализатор траектории → исполнитель (HM10 → RS-485) → замок")

# --- Слайд 6: алгоритм ---
s = blank()
title(s, "Алгоритм анализа траектории")
bullets(s, [
    "Конвейер: RSSI → фильтр Калмана → оценка дистанции → тренд (МНК) → конечный автомат",
    "Доступ выдаётся при устойчивом приближении носителя в зону",
], top=1.7, height=1.4, size=18)
image_fit(s, "fig_2_2_fsm.png", top=3.2, max_h=3.9,
          caption="Граф состояний принятия решения о доступе")

# --- Слайд 7: реализация ---
s = blank()
title(s, "Программная реализация")
bullets(s, [
    "Сканер и генератор BLE-меток: Python (bleak), Flutter (Android), iOS",
    "Анализатор траектории — модуль trajectory.py",
    "Мобильное приложение + модуль HM10 (мост BLE → RS-485), команда 10 байт",
    "Идентификатор носителя — стабильный токен (не MAC) или динамический rolling-code (HMAC-SHA256)",
    "Доп. канал доступа — по входящему звонку (номер в BCD)",
    "Виртуальная база замков для воспроизводимого тестирования",
])

# --- Слайд 8: тестирование (моделирование) ---
s = blank()
title(s, "Тестирование: имитационное моделирование")
image_fit(s, "fig_3_4_sim_result.png", top=2.0, max_h=4.4,
          caption="При приближении носителя доступ предоставляется, при удалении — сбрасывается")

# --- Слайд 9: порог vs траектория ---
s = blank()
title(s, "Преимущество перед пороговым подходом")
image_fit(s, "fig_2_3_threshold_vs_trajectory.png", top=2.0, max_h=4.2,
          caption="Статичная метка: порог даёт ложные срабатывания, анализ траектории — нет")

# --- Слайд 10: заключение ---
s = blank()
title(s, "Заключение")
bullets(s, [
    "Цель достигнута, все задачи решены",
    "Новизна: анализ траектории (Калман + дистанция + тренд + автомат) вместо порога",
    "Практическая значимость: программный комплекс для построения СКУД на BLE",
    "Защищённость: динамический идентификатор (rolling-code) и журналирование; доп. канал — по звонку",
    "Тестирование подтвердило работоспособность и устойчивость к ложным срабатываниям",
    "Развитие: несколько приёмников, машинное обучение, протокол «запрос–ответ» против relay",
])

# --- Слайд 11: спасибо ---
s = blank()
tb = s.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.5))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Спасибо за внимание!"
r.font.size = Pt(40)
r.font.bold = True
r.font.name = "Arial"
r.font.color.rgb = NAVY

prs.save(OUT)
print(f"Сохранено: {OUT}; слайдов: {len(prs.slides._sldIdLst)}")
